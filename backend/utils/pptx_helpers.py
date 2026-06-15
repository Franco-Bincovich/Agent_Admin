import io

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


def set_slide_background(slide, hex_color: str) -> None:
    """Sets a solid background color on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def add_textbox(
    slide, pos: dict, text: str, font: str, size: int, color: str,
    bold: bool = False, auto_size: bool = False
) -> None:
    """Adds a formatted text box at the given position dict.

    When auto_size=True the font shrinks to fit if the wrapped text
    exceeds the box height (MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE).
    """
    txb = slide.shapes.add_textbox(pos["left"], pos["top"], pos["width"], pos["height"])
    tf = txb.text_frame
    tf.word_wrap = True
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        run = p.runs[0]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.bold = bold


def add_rect(slide, left, top, width, height, color: str):
    """Adds a solid filled rectangle with no border. Returns the shape."""
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(color)
    sh.line.fill.background()
    return sh


def set_shape_text(sh, text: str, font: str, size: int, color: str, bold: bool = False) -> None:
    """Writes formatted text into an existing shape's text frame."""
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        run = p.runs[0]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.bold = bold


def _rrect(slide, left, top, width, height, color: str, radius: float = 0.05):
    """Rectángulo redondeado con fill sólido y sin borde. radius=0.5 da círculo perfecto."""
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(color)
    sh.line.fill.background()
    sh.adjustments[0] = radius
    return sh


def _titulo_font_size(titulo: str, default: int = 28) -> int:
    """Reduce el font size del título según longitud para evitar desborde."""
    if len(titulo) <= 35:
        return default
    elif len(titulo) <= 55:
        return max(22, default - 4)
    else:
        return max(18, default - 8)


def _titulo_font_size_portada(titulo: str) -> int:
    """Font size para portada y cierre: 48 / 36 / 28 según longitud."""
    if len(titulo) <= 25:
        return 48
    elif len(titulo) <= 40:
        return 36
    else:
        return 28


def _add_textbox_italic(slide, pos: dict, text: str, font: str, size: int, color: str) -> None:
    """Textbox con formato italic."""
    txb = slide.shapes.add_textbox(pos["left"], pos["top"], pos["width"], pos["height"])
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        run = p.runs[0]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.italic = True


def _insert_image(slide, imagen) -> None:
    """Inserta imagen a la derecha del contenido. Falla silenciosamente."""
    try:
        img_bytes = imagen if isinstance(imagen, bytes) else imagen[1]
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            left=int(Inches(7.8)), top=int(Inches(1.3)),
            width=int(Inches(5.0)), height=int(Inches(5.8)),
        )
    except Exception:
        pass
