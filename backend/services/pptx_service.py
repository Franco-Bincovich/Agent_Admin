from __future__ import annotations

import io
from typing import Optional

from pptx import Presentation
from pptx.util import Cm, Inches

from services.pptx_layouts import BUILDERS
from templates import corporativo_neutro, ejecutivo_oscuro, profesional_claro
from utils.errors import AppError, ErrorCode

TEMPLATES = {
    "ejecutivo_oscuro": ejecutivo_oscuro,
    "profesional_claro": profesional_claro,
    "corporativo_neutro": corporativo_neutro,
}


def _insert_logo_pptx(slide, logo_bytes: bytes) -> None:
    """Inserta logo pequeño en la esquina superior izquierda del slide. Falla silenciosamente."""
    try:
        slide.shapes.add_picture(io.BytesIO(logo_bytes), left=Cm(0.5), top=Cm(0.3), height=Cm(2.1))
    except Exception:
        pass


def generate_pptx(
    outline: dict,
    template_name: str,
    logo_bytes: Optional[bytes] = None,
    imagenes: Optional[list[bytes]] = None,
) -> bytes:
    """
    Genera un archivo .pptx a partir de un outline dict y el nombre de un template.

    Args:
        outline: dict con 'titulo_presentacion' (str) y 'slides' (list[dict]).
                 Cada slide tiene 'tipo', 'titulo' y 'contenido'.
        template_name: nombre del template. Valores: 'ejecutivo_oscuro' |
                       'profesional_claro' | 'corporativo_neutro'.
        logo_bytes: Bytes del logo a insertar en la portada (slide 0). None omite el logo.
        imagenes: Bytes de imágenes extraídas del documento fuente. Se insertan una por
                  slide en slides de tipo 'contenido' y 'destacado', en orden de
                  disponibilidad. None o lista vacía mantienen el comportamiento original.

    Returns:
        Bytes del archivo .pptx generado, listos para enviar como descarga.

    Raises:
        AppError: code TEMPLATE_NOT_FOUND (400) si template_name no existe en TEMPLATES.
        AppError: code PPTX_BUILD_ERROR (503) si la generación del archivo falla.
    """
    if template_name not in TEMPLATES:
        raise AppError("Template no encontrado", ErrorCode.TEMPLATE_NOT_FOUND, 400)
    tpl = TEMPLATES[template_name]
    _img_slides = {"contenido", "destacado"}
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        img_map: dict[int, bytes] = {}
        if imagenes:
            for i, img in enumerate(imagenes):
                img_map[i] = img
        for idx, slide_data in enumerate(outline["slides"]):
            slide = prs.slides.add_slide(blank_layout)
            builder = BUILDERS.get(slide_data["tipo"])
            if builder:
                if slide_data["tipo"] in _img_slides:
                    imagen_para_slide = None
                    img_idx = slide_data.get("imagen_idx")
                    if isinstance(img_idx, int) and img_idx in img_map:
                        imagen_para_slide = img_map[img_idx]
                    builder(slide, slide_data, tpl, imagen_para_slide)
                else:
                    builder(slide, slide_data, tpl)
            if idx == 0 and logo_bytes:
                _insert_logo_pptx(slide, logo_bytes)
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()
    except AppError:
        raise
    except Exception as exc:
        raise AppError("Error generando PPTX", ErrorCode.PPTX_BUILD_ERROR, 503) from exc
