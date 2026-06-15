import io
import zipfile
from pathlib import Path

import fitz
import mammoth
import openpyxl
from pptx import Presentation

from utils.errors import AppError, ErrorCode
from utils.logger import log

MIN_TEXT_LENGTH = 50


def _extract_from_pdf(file_bytes: bytes) -> str:
    """Extrae texto de un PDF página por página usando PyMuPDF."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(pages)


def _extract_from_docx(file_bytes: bytes) -> str:
    """Extrae texto plano de un DOCX conservando la estructura de párrafos."""
    result = mammoth.extract_raw_text(io.BytesIO(file_bytes))
    return result.value


def _extract_from_txt(file_bytes: bytes) -> str:
    """Decodifica un archivo TXT con fallback de utf-8 a latin-1."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


def _extract_from_pptx(file_bytes: bytes) -> str:
    """Extrae texto de un PPTX iterando slides en orden, separando título del resto."""
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title = text
            else:
                body_parts.append(text)
        header = f"## Slide {i}: {title}" if title else f"## Slide {i}"
        slides_text.append(header + ("\n" + "\n".join(body_parts) if body_parts else ""))
    return "\n\n".join(slides_text)


def _extract_from_xlsx(file_bytes: bytes) -> str:
    """Extrae texto de todas las hojas de un XLSX, celda por celda."""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                rows.append("\t".join(cells))
    return "\n".join(rows)


_EXTRACTORS = {
    ".pdf": _extract_from_pdf,
    ".docx": _extract_from_docx,
    ".txt": _extract_from_txt,
    ".xlsx": _extract_from_xlsx,
    ".pptx": _extract_from_pptx,
}


def _docx_has_embedded_media(file_bytes: bytes) -> bool:
    """
    Indica si un DOCX trae imágenes o EMF/WMF embebidos en word/media/.

    Sirve para distinguir un DOCX EMF-heavy (todo el contenido en imágenes,
    texto plano escaso) de un archivo realmente vacío. El primero debe seguir
    al pipeline para entrar por el canal de visión; el segundo no.

    Returns:
        True si existe al menos un recurso bajo word/media/. False si el ZIP
        no tiene media o no puede abrirse.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            return any(name.startswith("word/media/") for name in zf.namelist())
    except Exception:
        return False


def extract_text_from_file(filename: str, file_bytes: bytes) -> str:
    """
    Extrae texto plano de un archivo, detectando el tipo por extensión.

    Delega a la función específica según la extensión y valida que el
    resultado contenga contenido útil (mínimo 50 caracteres).

    Args:
        filename: Nombre del archivo con su extensión (ej. 'informe.pdf').
        file_bytes: Contenido binario del archivo.

    Returns:
        Texto extraído, limpio y sin espacios extremos.

    Raises:
        AppError: code 'UNSUPPORTED_FORMAT', status 400 si la extensión no está soportada.
        AppError: code 'EMPTY_FILE', status 422 si el texto tiene menos de 50 caracteres,
            salvo que sea un .docx con media embebida (EMF-heavy): en ese caso se deja
            pasar para que el canal de visión procese las imágenes.
    """
    ext = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if not extractor:
        raise AppError(
            f"Formato no soportado: '{ext}'. Formatos válidos: pdf, docx, txt, xlsx, pptx.",
            ErrorCode.UNSUPPORTED_FORMAT,
            400,
        )
    text = extractor(file_bytes)
    if len(text.strip()) < MIN_TEXT_LENGTH:
        if ext == ".docx" and _docx_has_embedded_media(file_bytes):
            log.info(
                "extract.docx_emf_heavy | texto plano escaso, sigue por "
                f"canal de visión | filename={filename} | chars={len(text.strip())}"
            )
            return text.strip()
        raise AppError(
            "El archivo no contiene texto suficiente.",
            ErrorCode.EMPTY_FILE,
            422,
        )
    return text.strip()
