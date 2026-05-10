from __future__ import annotations

import io

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Cm, Pt

from utils.pptx_helpers import add_textbox, set_slide_background

# Ancho estándar de una presentación PowerPoint 16:9 en EMU (English Metric Units).
_SLIDE_WIDTH_EMU = 9144000


def _build_portada(slide, data: dict, tpl) -> None:
    """
    Construye slide de portada: fondo sólido, título grande en accent_color,
    subtítulo debajo en secondary_text.

    Args:
        slide: objeto slide de python-pptx.
        data: dict con 'titulo' (str) y 'contenido' (str).
        tpl: módulo de template con atributos TEMPLATE y LAYOUTS.
    """
    template = tpl.TEMPLATE
    layout = tpl.LAYOUTS["portada"]
    set_slide_background(slide, template["background_color"])
    add_textbox(slide, layout["title"], data["titulo"], template["font_title"], template["font_size_title_portada"], template["accent_color"], bold=True)
    add_textbox(slide, layout["subtitle"], str(data["contenido"]), template["font_body"], template["font_size_body"], template["secondary_text"])


def _build_contenido(slide, data: dict, tpl, imagen: tuple[str, bytes] | None = None) -> None:
    """
    Construye slide de contenido: fondo, título en accent_color, bullets en text_color.
    Cada bullet lleva prefijo '•'. Se renderizan máximo 5. Si imagen está disponible,
    se inserta a la derecha del texto (ancho máximo 40% del slide).

    Args:
        slide: objeto slide de python-pptx.
        data: dict con 'titulo' (str) y 'contenido' (list[str]).
        tpl: módulo de template con atributos TEMPLATE y LAYOUTS.
        imagen: Bytes de imagen a insertar a la derecha, o None para omitir.
    """
    template = tpl.TEMPLATE
    layout = tpl.LAYOUTS["contenido"]
    set_slide_background(slide, template["background_color"])
    add_textbox(slide, layout["title"], data["titulo"], template["font_title"], template["font_size_title_slide"], template["accent_color"], bold=True)
    bullets = data["contenido"] if isinstance(data["contenido"], list) else [str(data["contenido"])]
    bullet_text = "\n".join(f"• {b}" for b in bullets[:5])
    add_textbox(slide, layout["body"], bullet_text, template["font_body"], template["font_size_body"], template["text_color"])
    if imagen:
        _insert_image_on_slide(slide, imagen)


def _build_destacado(slide, data: dict, tpl, imagen: tuple[str, bytes] | None = None) -> None:
    """
    Construye slide destacado: fondo, título en accent_color, texto central en
    rectángulo relleno (fondo accent_color, texto background_color). Si imagen está
    disponible, se inserta a la derecha del texto (ancho máximo 40% del slide).

    Args:
        slide: objeto slide de python-pptx.
        data: dict con 'titulo' (str) y 'contenido' (str).
        tpl: módulo de template con atributos TEMPLATE y LAYOUTS.
        imagen: Bytes de imagen a insertar a la derecha, o None para omitir.
    """
    template = tpl.TEMPLATE
    layout = tpl.LAYOUTS["destacado"]
    set_slide_background(slide, template["background_color"])
    add_textbox(slide, layout["title"], data["titulo"], template["font_title"], template["font_size_title_slide"], template["accent_color"], bold=True)
    box = layout["box"]
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, box["left"], box["top"], box["width"], box["height"])
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(template["accent_color"])
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = str(data["contenido"])
    run = tf.paragraphs[0].runs[0]
    run.font.name = template["font_body"]
    run.font.size = Pt(template["font_size_body"])
    run.font.color.rgb = RGBColor.from_string(template["background_color"])
    if imagen:
        _insert_image_on_slide(slide, imagen)


def _build_cierre(slide, data: dict, tpl) -> None:
    """
    Construye slide de cierre: fondo, texto centrado en accent_color,
    línea horizontal decorativa debajo en secondary_text.

    Args:
        slide: objeto slide de python-pptx.
        data: dict con 'titulo' (str) y 'contenido' (str).
        tpl: módulo de template con atributos TEMPLATE y LAYOUTS.
    """
    template = tpl.TEMPLATE
    layout = tpl.LAYOUTS["cierre"]
    set_slide_background(slide, template["background_color"])
    add_textbox(slide, layout["text"], data["titulo"], template["font_title"], template["font_size_title_portada"], template["accent_color"], bold=True)
    ln = layout["line"]
    line_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, ln["left"], ln["top"], ln["width"], ln["height"])
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor.from_string(template["secondary_text"])
    line_shape.line.fill.background()


def _insert_image_on_slide(slide, imagen: tuple[str, bytes]) -> None:
    """
    Inserta una imagen a la derecha del texto en el slide. Falla silenciosamente.

    La imagen ocupa un ancho máximo del 40% del slide estándar (9144000 EMU) y se
    posiciona en el margen derecho, comenzando a 2.5 cm del borde superior.
    """
    try:
        _, img_bytes = imagen
        max_width = int(_SLIDE_WIDTH_EMU * 0.4)
        left = _SLIDE_WIDTH_EMU - max_width - int(Cm(0.5))
        top = int(Cm(2.5))
        slide.shapes.add_picture(io.BytesIO(img_bytes), left=left, top=top, width=max_width)
    except Exception:
        pass
